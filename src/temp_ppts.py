"""
实验脚本
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE

from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

from pptx.chart.xmlwriter import SeriesXmlRewriterFactory, ChartXmlWriter
from lxml import etree
from io import BytesIO
from pptx.oxml.ns import _nsmap as namespaces

data = {
    '5%': [370, 381],
    '25%':[435, 458],
    '50%-25%': [71, 89],
    '75%-50%': [68, 79],
    '95%': [673, 724]
}

# XMLs used to set the properties 
# 1. set serie fill to transparent 
xml_nf = '''<c:spPr %s><a:noFill/></c:spPr>'''% nsdecls("c", "a")
# 2. set line fill to none 
lnf_xml = '''<c:spPr %s>
    <a:ln w="25400">
        <a:noFill/>
    </a:ln>
    <a:effectLst/>
</c:spPr>''' %  nsdecls("c", "a")
# 3. Add a hihe-low line 
hi_low_xml = "<c:hiLowLines %s/>" % nsdecls("c") 

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
chart_data = CategoryChartData()
chart_data.categories = ['One', 'Two']
for d in ['25%', '50%-25%', '75%-50%']:
    chart_data.add_series(d, data[d])

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
column_bars = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data
)

# Hide the bar in the bottom showing 25 percentile 
ser_one_no_fill = parse_xml(xml_nf)
column_bars.chart._chartSpace.plotArea.xpath("*/c:ser")[0].insert(0, ser_one_no_fill)

# Add a line plot in chartspace/plotarea 
line_data = CategoryChartData()
line_data.categories = ['One', 'Two']
for d in ['5%', '95%']:
    line_data.add_series(d, data[d])

xml_bytes = ChartXmlWriter(XL_CHART_TYPE.LINE, line_data).xml
root_el = etree.parse(BytesIO(xml_bytes.encode('utf-8')))

val_axis = column_bars.chart._chartSpace.plotArea.xpath("c:valAx/c:axId/@val")
cat_axis = column_bars.chart._chartSpace.plotArea.xpath("c:catAx/c:axId/@val")
series_count = len(column_bars.chart._chartSpace.plotArea.xpath("*/c:ser"))

for el in root_el.xpath("/c:chartSpace/c:chart/c:plotArea/*", namespaces=namespaces):
    if el.tag.endswith("Chart"):  
        # Since the two charts share the same axies, the axId should be the same 
        existing_axis_ref = el.xpath("c:axId", namespaces=namespaces)
        for x in existing_axis_ref:
            el.remove(x)
        for x in cat_axis + val_axis:   # Order is important 
            el.append(etree.Element("{%s}axId" % namespaces['c'], val=x))

        # Need to increment series index to make PP happy
        for i, x in enumerate(el.xpath("c:ser/c:idx", namespaces=namespaces)):
            x.set("val", str(series_count + i))
        for i, x in enumerate(el.xpath("c:ser/c:order", namespaces=namespaces)):
            x.set("val", str(series_count + i))
        
        # Show a short dash at 5 and 95 percentiles 
        markers = el.xpath("c:ser/c:marker/c:symbol", namespaces=namespaces)
        for marker in markers:
            marker.set('val', 'dash')
        
        # Hide the line 
        for s in  el.xpath("c:ser", namespaces = namespaces):
            lnf = parse_xml(lnf_xml)
            s.append(lnf)
        
        # Add a high-low line 
        el.append(parse_xml(hi_low_xml))
        column_bars.chart._chartSpace.plotArea.insert(1, el)  

prs.save('output/test.pptx')