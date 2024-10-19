import time
start = time.time()
# mysetup = """
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt, Cm
import os
import pdb
# import polars as pl


## pptx report
# Title: IR1 异常值 wafer map 分类
# Chart: x: 时间， Y：IR1 ratio by wafer
# Table: 同一X 轴， Y: 各缺陷率占比
# Table：X: SSA Defect class, Y: row01-map and row02-image

"""PPTx struct
-/Presentation
--/slides
---/shapes
----/组件(title, text, chart, table, image...)
 MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
"""

class report():
    
    def __init__(self):
        self.prs = Presentation()

    def define_theme(self, data_df, conf):
        slide_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(slide_layout)
        shapes = self.slide.shapes

        # ======>>> 设置 title
        left = Cm(0.5)
        top = Cm(-0.8)    # left，top为相对位置
        width = Cm(24.5)    # width，height为文本框的大小
        height = Cm(1)
        textbox = shapes.add_textbox(left, top, width, height)
        title = textbox.text_frame.add_paragraph()
        # title.alignment = PP_ALIGN.LEFT
        title.text = conf["title"]
        title.font.name = 'Calibri'
        title.font.size = Pt(18)
        title.font.bold = True
        title.font.italic = True
        title.font.color.rgb = RGBColor(0x15, 0x21, 0x77)

        # ======>>> 设置图表
        left = Cm(0.5)
        top = Cm(1)
        width = Cm(24.5)
        height = Cm(5)
        chart_data = ChartData()
        chart_data.categories = ['Q1 Sales', 'Q2 Sales', 'Q3 Sales']
        chart_data.add_series('IR', (32.2, 58.4, 14.7))
        # chart_data.add_series('VF', (2.2, 28.4, 44.7))
         
        graphic_frame = shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS_STACKED, left, top, width, height, chart_data
            )
        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.series[0].smooth = True


        # ======>>> 设置 chart 表格
        rows, cols = data_df["chart_table_layout"]
        left = Cm(0.5)
        top = Cm(6)
        width = Cm(24.5)
        height = Cm(3)
        graphic_frame = shapes.add_table(rows, cols, left, top, width, height)

        IR_table = graphic_frame.table
        IR_table.first_col = False
        IR_table.first_row = False
        for col in range(cols):
            if 0 == col:
                IR_table.columns[col].width = Cm(6.5)
            else:
                IR_table.columns[col].width = Cm(1)
        for row in range(rows):

            IR_table.rows[row].height = Cm(0.25)
            IR_table.cell(row,0).text = list(data_df["abnormal_item"]["IR_ratio"].keys())[row]

        # ======>>> 设置Map表格
        rows, cols = data_df["map_table_layout"]
        left = Cm(0.5)
        top = Cm(11)
        width = Cm(24.5)
        height = Cm(8)
        graphic_frame = shapes.add_table(rows, cols, left, top, width, height)
        abnor_table = graphic_frame.table
        # abnor_table.ApplyStyle(2) 
        # 启用 表头特殊格式
        abnor_table.first_row = True
        abnor_table.horz_banding = False
        abnor_table.vert_banding = False
        abnor_table.columns[0].width = Cm(3)
        abnor_table.rows[0].height = Cm(0.5)

        for cel_h in range(cols):
            if cel_h > 0:
                abnor_table.columns[cel_h].width = Cm(5.35)
            # abnor_table.cell(0,cel_h).fill.fore_color.rgb = RGBColor(0x15, 0x21, 0x77) # 红色
            # print(f"=====>>> {abnor_table.cell(0,cel_h).fill.fore_color()}")
            cell_color = abnor_table.cell(0,cel_h)
            # cell_color.fill.background().solid().foregroundColor = RGBColor(0x15, 0x21, 0x77)
            cell = abnor_table.cell(0,cel_h)
            # for border in cell.borders:
            #     border.color.rgb = RGBColor(0x15, 0x21, 0x77)
            # print(f"=====>>> {cell_color}")
            # cell_color = cell_color.gradient()
            title = abnor_table.cell(0,cel_h).text_frame
            title.alignment = PP_ALIGN.CENTER
            title.text = data_df["abnormal_item"]["schema"]["column"][cel_h]

        for cel_v in range(rows):
            # breakpoint()
            if 0 == cel_v:
                abnor_table.rows[cel_v].height = Cm(0.5)
            if 1 == cel_v:
                abnor_table.rows[cel_v].height = Cm(4.5)
            if 2 == cel_v:
                abnor_table.rows[cel_v].height = Cm(2.5)
            if cel_v > 0 and cel_v < rows:
                abnor_table.cell(cel_v,0).text = data_df["abnormal_item"]["schema"]["left_row"][cel_v-1]
                abnor_table.cell(cel_v,0).alignment = PP_ALIGN.CENTER

            else:
                pass

        # ====== 设置map图片
        left = Cm(3.5+0.5)
        top = Cm(12)
        width = height = Cm(4.5)
        map_pics = data_df["abnormal_item"]["mapping"]
        for pic in map_pics.keys():
            shapes.add_picture(map_pics[pic], left, top, width, height)
            left += width + Cm(0.8)

        
        # ====== 设置 OM 图片
        left = Cm(3.5+1)
        top = Cm(16.5)
        width = height = Cm(2.5)
        map_pics = data_df["abnormal_item"]["om"]
        
        for pic in map_pics.keys():
            shapes.add_picture(map_pics[pic], left, top, width, height)
            left += width + Cm(3) 
        self.prs.save(conf["output_dir"])
        self.prs.save("output/pptx-tmp.xml")
        # self.slide.saveas("output/pptx-tmp.png", format='png')



    def save_as(self, img_format, output):
        pass

if __name__ == "__main__":
    r = report()
    data_df = {
        "chart_table_layout" : (5,6),
        "map_table_layout" : (3,5),
        "abnormal_item": {
        "title":"异常项",
        "IR_ratio" : {"石磨盘破损占比":{},
                "区域性异常占比":{},
                "疑似划伤或污染占比":{},
                "外圈或散点分布占比":{},
                "特殊异常占比":{},},
        "schema":{
            "column" :["异常项","石磨盘破损","区域性异常","疑似划伤或污染","外圈或散点分布"],
            "left_row":["Mapping","OM"],
            },
        "mapping" : {"area":os.path.abspath(r"test/map/NA057230608C01.png"),
                   "circle":os.path.abspath(r"test/map/NA057230608C01.png"),
                   "scratch_dirty":os.path.abspath(r"test/map/NA057230608C01.png"),
                   "smp":os.path.abspath(r"test/map/NA057230608C01.png"),
                    },
        "om" : {"area":os.path.abspath(r"test/om/area.bmp"),
                "circle":os.path.abspath(r"test/om/circle.bmp"),
                "scratch_dirty":os.path.abspath(r"test/om/scratch_dirty.bmp"),
                "smp":os.path.abspath(r"test/om/smp.bmp"),
                },
            },
        "IR_chart" :[],
        }
    temp_conf = {
        "slides" : 1,
        "output_dir":"output/chart-03.pptx",
        "title" : "IR1 < 90% wafer map 分类报告",
        "time_scope" : "3months",
    }

    r.define_theme(data_df, temp_conf)
    end = time.time()
    print(f"Iteration: start time: {start}, end time: {end}, taken: {(end-start)} s")

